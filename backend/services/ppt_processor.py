"""
services/ppt_processor.py
Extract text from .pptx / .ppt files using python-pptx.
"""
import os


def process_pptx(file_path: str, file_id: str, filename: str,
                 session_id: str = None, project_id: str = None) -> dict:
    """
    Extract all text from a PowerPoint file.
    Returns { success, text, preview, slides, filename }
    """
    if not os.path.exists(file_path):
        return {"success": False, "error": "File not found"}

    try:
        from pptx import Presentation
    except ImportError:
        return {
            "success": False,
            "error": "python-pptx not installed.\nRun: pip install python-pptx --break-system-packages"
        }

    try:
        prs    = Presentation(file_path)
        slides = []

        for i, slide in enumerate(prs.slides, 1):
            slide_parts = []

            # Slide title
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_parts.append(f"[Title] {slide.shapes.title.text.strip()}")

            # All text from shapes
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if shape == slide.shapes.title:
                    continue  # already added
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)

            # Tables
            for shape in slide.shapes:
                if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                    try:
                        tbl = shape.table
                        rows = []
                        for row in tbl.rows:
                            cells = [c.text.strip() for c in row.cells]
                            rows.append(" | ".join(cells))
                        slide_parts.append("\n".join(rows))
                    except Exception:
                        pass

            # Speaker notes
            try:
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_parts.append(f"[Notes] {notes_text}")
            except Exception:
                pass

            if slide_parts:
                slides.append(f"--- Slide {i} ---\n" + "\n".join(slide_parts))

        full_text = "\n\n".join(slides)

        # Store in vector memory for RAG
        try:
            from services.vector_memory import store_document
            ns = project_id or session_id or file_id
            store_document(file_id, full_text, ns, filename, "pptx")
        except Exception:
            pass

        return {
            "success":  True,
            "file_id":  file_id,
            "filename": filename,
            "slides":   len(slides),
            "chars":    len(full_text),
            "text":     full_text,
            "preview":  full_text[:500],
        }

    except Exception as e:
        return {"success": False, "error": f"Failed to read PPTX: {str(e)}"}
